from fastapi import FastAPI, APIRouter, HTTPException, Depends, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.responses import FileResponse, StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict, EmailStr
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone, timedelta
import bcrypt
import jwt

from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
import io
import asyncio
# import google.generativeai as genai
import google.generativeai as genai

class LlmChat:
    def __init__(self, api_key: str, session_id: str, system_message: str):
        genai.configure(api_key=api_key)
        self.system_message = system_message
        self.session_id = session_id
        self.model = None

    def with_model(self, provider: str, model_name: str):
        self.model = genai.GenerativeModel(model_name)
        return self

    async def send_message(self, user_message):

        loop = asyncio.get_event_loop()

        def _call():

            return self.model.generate_content([
            {
                "role": "model",
                "parts": [{"text": self.system_message}]
            },
            {
                "role": "user",
                "parts": [{"text": user_message.text}]
            }
        ])

        response = await loop.run_in_executor(None, _call)
        return response.text





class UserMessage:
    def __init__(self, text: str):
        self.text = text



ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
# client = AsyncIOMotorClient(mongo_url)
client = AsyncIOMotorClient(
    mongo_url,
    tlsAllowInvalidCertificates=True
)

db = client[os.environ['DB_NAME']]

# JWT Configuration
JWT_SECRET_KEY = os.environ.get('JWT_SECRET_KEY')
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_HOURS = 24

# LLM Configuration
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')

app = FastAPI()
api_router = APIRouter(prefix="/api")
security = HTTPBearer()

# ============== Models ==============

class UserRegister(BaseModel):
    email: EmailStr
    password: str
    name: str

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class User(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    email: str
    name: str
    created_at: str

class TokenResponse(BaseModel):
    token: str
    user: User

class ProjectCreate(BaseModel):
    title: str
    document_type: str  # "docx" or "pptx"
    topic: str

class ProjectStructure(BaseModel):
    sections: Optional[List[Dict[str, Any]]] = None  # For Word: [{"id": "", "title": "", "order": 0}]
    slides: Optional[List[Dict[str, Any]]] = None     # For PPT: [{"id": "", "title": "", "order": 0}]

class Project(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    user_id: str
    title: str
    document_type: str
    topic: str
    structure: Optional[ProjectStructure] = None
    status: str  # "draft", "configured", "generating", "generated", "completed"
    created_at: str
    updated_at: str

class ContentItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    project_id: str
    section_id: str  # or slide_id
    content: str
    version: int
    created_at: str

class RefinementRequest(BaseModel):
    section_id: str
    prompt: str

class FeedbackRequest(BaseModel):
    section_id: str
    feedback_type: str  # "like" or "dislike"
    comment: Optional[str] = None

class AIOutlineRequest(BaseModel):
    topic: str
    document_type: str

# ============== Auth Utilities ==============

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str) -> str:
    expiration = datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    payload = {
        "user_id": user_id,
        "exp": expiration
    }
    return jwt.encode(payload, JWT_SECRET_KEY, algorithm=JWT_ALGORITHM)

def decode_token(token: str) -> str:
    try:
        payload = jwt.decode(token, JWT_SECRET_KEY, algorithms=[JWT_ALGORITHM])
        return payload["user_id"]
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)) -> str:
    token = credentials.credentials
    return decode_token(token)

# ============== Auth Endpoints ==============

@api_router.post("/auth/register", response_model=TokenResponse)
async def register(user_data: UserRegister):
    # Check if user exists
    existing_user = await db.users.find_one({"email": user_data.email})
    if existing_user:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create user
    user_id = str(uuid.uuid4())
    user_doc = {
        "id": user_id,
        "email": user_data.email,
        "password_hash": hash_password(user_data.password),
        "name": user_data.name,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.users.insert_one(user_doc)
    
    # Create token
    token = create_token(user_id)
    user = User(
        id=user_id,
        email=user_data.email,
        name=user_data.name,
        created_at=user_doc["created_at"]
    )
    
    return TokenResponse(token=token, user=user)

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(credentials: UserLogin):
    user_doc = await db.users.find_one({"email": credentials.email})
    if not user_doc or not verify_password(credentials.password, user_doc["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    
    token = create_token(user_doc["id"])
    user = User(
        id=user_doc["id"],
        email=user_doc["email"],
        name=user_doc["name"],
        created_at=user_doc["created_at"]
    )
    
    return TokenResponse(token=token, user=user)

@api_router.get("/auth/me", response_model=User)
async def get_me(user_id: str = Depends(get_current_user)):
    user_doc = await db.users.find_one({"id": user_id})
    if not user_doc:
        raise HTTPException(status_code=404, detail="User not found")
    
    return User(
        id=user_doc["id"],
        email=user_doc["email"],
        name=user_doc["name"],
        created_at=user_doc["created_at"]
    )

# ============== Project Endpoints ==============

@api_router.post("/projects", response_model=Project)
async def create_project(project_data: ProjectCreate, user_id: str = Depends(get_current_user)):
    project_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    project_doc = {
        "id": project_id,
        "user_id": user_id,
        "title": project_data.title,
        "document_type": project_data.document_type,
        "topic": project_data.topic,
        "structure": None,
        "status": "draft",
        "created_at": now,
        "updated_at": now
    }
    await db.projects.insert_one(project_doc)
    
    return Project(**project_doc)

@api_router.get("/projects", response_model=List[Project])
async def get_projects(user_id: str = Depends(get_current_user)):
    projects = await db.projects.find({"user_id": user_id}, {"_id": 0}).to_list(1000)
    return [Project(**p) for p in projects]

@api_router.get("/projects/{project_id}", response_model=Project)
async def get_project(project_id: str, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    return Project(**project)

@api_router.put("/projects/{project_id}/structure")
async def update_structure(project_id: str, structure: ProjectStructure, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {
            "structure": structure.model_dump(),
            "status": "configured",
            "updated_at": datetime.now(timezone.utc).isoformat()
        }}
    )
    
    return {"message": "Structure updated successfully"}

@api_router.delete("/projects/{project_id}")
async def delete_project(project_id: str, user_id: str = Depends(get_current_user)):
    result = await db.projects.delete_one({"id": project_id, "user_id": user_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Delete associated content and feedback
    await db.content.delete_many({"project_id": project_id})
    await db.feedback.delete_many({"project_id": project_id})
    
    return {"message": "Project deleted successfully"}

# ============== AI Generation ==============

@api_router.post("/projects/{project_id}/generate")
async def generate_content(project_id: str, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if not project.get("structure"):
        raise HTTPException(status_code=400, detail="Project structure not configured")
    
    # Update status to generating
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {"status": "generating", "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    # Initialize LLM chat
    chat = LlmChat(
        api_key=GEMINI_API_KEY,
        session_id=f"project_{project_id}",
        system_message=f"You are an expert document writer. Generate professional, well-structured content for {project['document_type']} documents."
    ).with_model("gemini", "gemini-2.5-flash")
    
    structure = project["structure"]
    items = structure.get("sections") if project["document_type"] == "docx" else structure.get("slides")
    
    # Generate content for each section/slide
    for item in items:
        section_id = item["id"]
        title = item["title"]
        
        # Check if content already exists
        existing_content = await db.content.find_one({"project_id": project_id, "section_id": section_id})
        if existing_content:
            continue
        
        # Generate content
        prompt = f"""Topic: {project['topic']}
Section/Slide Title: {title}

Generate detailed, professional content for this section. 
{'If this is a slide, keep it concise with bullet points.' if project['document_type'] == 'pptx' else 'Provide comprehensive paragraphs with good structure.'}
Provide only the content, without any preamble or meta-commentary."""
        
        user_message = UserMessage(text=prompt)
        response = await chat.send_message(user_message)
        
        # Save content
        content_id = str(uuid.uuid4())
        content_doc = {
            "id": content_id,
            "project_id": project_id,
            "section_id": section_id,
            "content": response,
            "version": 1,
            "created_at": datetime.now(timezone.utc).isoformat()
        }
        await db.content.insert_one(content_doc)
    
    # Update status to generated
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {"status": "generated", "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    
    return {"message": "Content generated successfully"}

@api_router.get("/projects/{project_id}/content")
async def get_content(project_id: str, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all content for this project
    content_list = await db.content.find({"project_id": project_id}, {"_id": 0}).to_list(1000)
    
    # Get latest version for each section
    section_content = {}
    for content in content_list:
        section_id = content["section_id"]
        if section_id not in section_content or content["version"] > section_content[section_id]["version"]:
            section_content[section_id] = content
    
    return {"content": list(section_content.values())}

@api_router.post("/projects/{project_id}/refine")
async def refine_content(project_id: str, refinement: RefinementRequest, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get current content
    current_content = await db.content.find_one(
        {"project_id": project_id, "section_id": refinement.section_id},
        sort=[("version", -1)]
    )
    if not current_content:
        raise HTTPException(status_code=404, detail="Content not found")
    
    # Initialize LLM chat
    chat = LlmChat(
        api_key=GEMINI_API_KEY,
        session_id=f"refine_{project_id}_{refinement.section_id}",
        system_message="You are an expert content editor. Refine the provided content based on user feedback."
    ).with_model("gemini", "gemini-2.5-flash")
    
    # Refine content
    prompt = f"""Current content:
{current_content['content']}

User refinement request: {refinement.prompt}

Provide the refined content. Only return the improved content, without any preamble or explanation."""
    
    user_message = UserMessage(text=prompt)
    refined_content = await chat.send_message(user_message)
    
    # Save refined content as new version
    content_id = str(uuid.uuid4())
    new_version = current_content["version"] + 1
    content_doc = {
        "id": content_id,
        "project_id": project_id,
        "section_id": refinement.section_id,
        "content": refined_content,
        "version": new_version,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.content.insert_one(content_doc)
    
    return {"message": "Content refined successfully", "content": refined_content, "version": new_version}

@api_router.post("/projects/{project_id}/feedback")
async def submit_feedback(project_id: str, feedback: FeedbackRequest, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    feedback_id = str(uuid.uuid4())
    feedback_doc = {
        "id": feedback_id,
        "project_id": project_id,
        "section_id": feedback.section_id,
        "feedback_type": feedback.feedback_type,
        "comment": feedback.comment,
        "timestamp": datetime.now(timezone.utc).isoformat()
    }
    await db.feedback.insert_one(feedback_doc)
    
    return {"message": "Feedback submitted successfully"}

# ============== AI Outline Suggestion ==============

@api_router.post("/ai/suggest-outline")
async def suggest_outline(request: AIOutlineRequest, user_id: str = Depends(get_current_user)):
    chat = LlmChat(
        api_key=GEMINI_API_KEY,
        session_id=f"outline_{user_id}_{uuid.uuid4()}",
        system_message="You are an expert document planner. Generate structured outlines for documents."
    ).with_model("gemini", "gemini-2.5-flash")
    
    if request.document_type == "docx":
        prompt = f"""Generate a detailed outline for a Word document about: {request.topic}

Provide 5-7 section titles that would make a comprehensive document. Return only the section titles, one per line, without numbering or bullets."""
    else:
        prompt = f"""Generate a slide deck outline for a PowerPoint presentation about: {request.topic}

Provide 6-10 slide titles that would make an engaging presentation. Return only the slide titles, one per line, without numbering or bullets."""
    
    user_message = UserMessage(text=prompt)
    response = await chat.send_message(user_message)
    
    # Parse response into list
    titles = [line.strip() for line in response.strip().split('\n') if line.strip()]
    
    # Create structure
    items = []
    for idx, title in enumerate(titles):
        items.append({
            "id": str(uuid.uuid4()),
            "title": title,
            "order": idx
        })
    
    result = {"sections": items} if request.document_type == "docx" else {"slides": items}
    return result

# ============== Document Export ==============

@api_router.get("/projects/{project_id}/export")
async def export_document(project_id: str, user_id: str = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "user_id": user_id})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all content
    content_list = await db.content.find({"project_id": project_id}, {"_id": 0}).to_list(1000)
    
    # Get latest version for each section
    section_content = {}
    for content in content_list:
        section_id = content["section_id"]
        if section_id not in section_content or content["version"] > section_content[section_id]["version"]:
            section_content[section_id] = content
    
    structure = project["structure"]
    items = structure.get("sections") if project["document_type"] == "docx" else structure.get("slides")
    
    if project["document_type"] == "docx":
        # Create Word document
        doc = Document()
        
        # Add title
        title = doc.add_heading(project["title"], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add topic
        topic_para = doc.add_paragraph(f"Topic: {project['topic']}")
        topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()  # Empty line
        
        # Add sections
        for item in sorted(items, key=lambda x: x["order"]):
            section_id = item["id"]
            if section_id in section_content:
                # Add section heading
                doc.add_heading(item["title"], 1)
                
                # Add content
                content = section_content[section_id]["content"]
                doc.add_paragraph(content)
                doc.add_paragraph()  # Empty line
        
        # Save to bytes
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        filename = f"{project['title'].replace(' ', '_')}.docx"
        
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    else:
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = project["title"]
        subtitle.text = project["topic"]
        
        # Content slides
        for item in sorted(items, key=lambda x: x["order"]):
            section_id = item["id"]
            if section_id in section_content:
                # Use title and content layout
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = item["title"]
                
                # Set content
                content_box = slide.placeholders[1]
                content = section_content[section_id]["content"]
                content_box.text = content
        
        # Save to bytes
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        filename = f"{project['title'].replace(' ', '_')}.pptx"
        
        return StreamingResponse(
            file_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )

# ============== Health Check ==============

@api_router.get("/")
async def root():
    return {"message": "DocuCraft AI API is running"}

@api_router.get("/health")
async def health_check():
    return {"status": "healthy"}

# Include router
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()